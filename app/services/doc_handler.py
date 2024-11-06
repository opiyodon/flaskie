from pdfminer.high_level import extract_text as pdf_extract_text
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from PIL import Image
import pytesseract
import os

class DocHandler:
    def analyze_document(self, filepath):
        """Analyze a document and return relevant information."""
        file_ext = filepath.split('.')[-1].lower()
        
        if file_ext == 'pdf':
            return self._analyze_pdf(filepath)
        elif file_ext in ['doc', 'docx']:
            return self._analyze_docx(filepath)
        elif file_ext in ['xls', 'xlsx']:
            return self._analyze_excel(filepath)
        elif file_ext in ['ppt', 'pptx']:
            return self._analyze_powerpoint(filepath)
        elif file_ext in ['png', 'jpg', 'jpeg']:
            return self._analyze_image(filepath)
        else:
            raise ValueError(f'Unsupported file type: {file_ext}')

    def convert_document(self, filepath, target_format):
        """Convert a document to the target format."""
        # Implementation would depend on specific conversion needs
        pass

    def extract_text(self, filepath):
        """Extract text from various document types."""
        file_ext = filepath.split('.')[-1].lower()
        
        if file_ext == 'pdf':
            return pdf_extract_text(filepath)
        elif file_ext in ['doc', 'docx']:
            return self._extract_text_from_docx(filepath)
        elif file_ext in ['xls', 'xlsx']:
            return self._extract_text_from_excel(filepath)
        elif file_ext in ['ppt', 'pptx']:
            return self._extract_text_from_powerpoint(filepath)
        elif file_ext in ['png', 'jpg', 'jpeg']:
            return self._extract_text_from_image(filepath)
        else:
            raise ValueError(f'Unsupported file type: {file_ext}')

    def _analyze_pdf(self, filepath):
        """Analyze a PDF document."""
        text = pdf_extract_text(filepath)
        return {
            'type': 'pdf',
            'text_length': len(text),
            'pages': len(text.split('\f')),
            'preview': text[:500]
        }

    def _analyze_docx(self, filepath):
        """Analyze a Word document."""
        doc = Document(filepath)
        paragraphs = [p.text for p in doc.paragraphs]
        return {
            'type': 'docx',
            'paragraphs': len(paragraphs),
            'text_length': sum(len(p) for p in paragraphs),
            'preview': paragraphs[0] if paragraphs else ''
        }

    def _analyze_excel(self, filepath):
        """Analyze an Excel workbook."""
        wb = load_workbook(filepath, read_only=True)
        sheets_info = {}
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            sheets_info[sheet] = {
                'max_row': ws.max_row,
                'max_column': ws.max_column
            }
        return {
            'type': 'excel',
            'sheets': sheets_info
        }

    def _analyze_powerpoint(self, filepath):
        """Analyze a PowerPoint presentation."""
        prs = Presentation(filepath)
        return {
            'type': 'powerpoint',
            'slides': len(prs.slides),
            'layouts': len(prs.slide_layouts)
        }

    def _analyze_image(self, filepath):
        """Analyze an image."""
        with Image.open(filepath) as img:
            return {
                'type': 'image',
                'format': img.format,
                'mode': img.mode,
                'size': img.size
            }

    def _extract_text_from_docx(self, filepath):
        """Extract text from a Word document."""
        doc = Document(filepath)
        return '\n'.join(p.text for p in doc.paragraphs)

    def _extract_text_from_excel(self, filepath):
        """Extract text from an Excel workbook."""
        wb = load_workbook(filepath, read_only=True)
        text = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.rows:
                text.extend(str(cell.value) for cell in row if cell.value)
        return '\n'.join(text)

    def _extract_text_from_powerpoint(self, filepath):
        """Extract text from a PowerPoint presentation."""
        prs = Presentation(filepath)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text.append(shape.text)
        return '\n'.join(text)

    def _extract_text_from_image(self, filepath):
        """Extract text from an image using OCR."""
        return pytesseract.image_to_string(Image.open(filepath))
